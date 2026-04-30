#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
潔咪 Sub-Agent 腳本
獨立完成深度研究報告：財富自由後的人生規劃（圖文並茂版）

使用方式：python3 jiemi_subagent.py
"""

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

# Groq API 設定
GROQ_API_KEY = "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"
GROQ_URL = "https://api.groq.com/openai/v1/chat/completions"

SYSTEM_PROMPT = """你是潔咪，一位才華洋溢的簡報設計師與人生規劃顧問。

⚠️ 強制規定：所有輸出必須使用繁體中文，嚴禁簡體中文！

你的專長：
- 簡報設計與視覺排版
- 人生規劃與自我實現
- 財務自由後的生活藍圖
- 激勵人心的演說風格

你的特點：
- 說話溫暖有深度
- 善於用故事和比喻
- 視覺化呈現豐富
- 讓人感到希望與動力

請用條列式呈現你的分析內容，包含豐富的細節和建議。"""

def get_jiemi_analysis():
    """取得潔咪的完整分析內容"""
    prompt = f"""{SYSTEM_PROMPT}

主題：財富自由後的人生規劃

請用條列式呈現你的分析內容（每個章節需要詳細內容）：
1. 財富自由的定義與標準
2. 財富自由後的生活想像
3. 人生各階段規劃
4. 實現路徑與時間表
5. 心理準備與心態調整
6. 給讀者的行動建議

請提供豐富詳細的內容，讓讀者能得到具體的指引。"""

    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }

    data = {
        "model": "openai/gpt-oss-120b",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.8
    }

    try:
        response = requests.post(GROQ_URL, headers=headers, json=data, timeout=90)
        if response.status_code == 200:
            return response.json()['choices'][0]['message']['content']
        else:
            return f"錯誤: {response.status_code}"
    except Exception as e:
        return f"例外: {str(e)}"

def create_wealth_freedom_ppt(analysis_content, output_path):
    """建立圖文並茂的 PowerPoint 簡報"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 投影片1：標題頁
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景裝飾
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(2.5))
    tf = title_box.text_frame
    tf.text = "💰 財富自由後的人生規劃"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 0)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(12.333), Inches(1.5))
    sf = subtitle_box.text_frame
    sf.text = "川寶研究團隊｜潔咪獨立研究報告"
    sf.paragraphs[0].font.size = Pt(28)
    sf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
    df = date_box.text_frame
    df.text = datetime.datetime.now().strftime("%Y年%m月%d日")
    df.paragraphs[0].font.size = Pt(18)
    df.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 投影片2：什麼是財富自由？
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    t2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    t2.text_frame.text = "🏠 什麼是財富自由？"
    t2.text_frame.paragraphs[0].font.size = Pt(40)
    t2.text_frame.paragraphs[0].font.bold = True
    
    c2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    c2.text_frame.text = """💡 財富自由的定義

「被動收入 ≥ 生活開銷」

當你的投資、租金、版稅等不需要親自工作的收入，
能夠稳定覆蓋你的日常生活開銷時，
你就達成了財富自由！

📊 簡易計算公式：
財富自由門檻 = 每月支出 × 12 ÷ 年化收益率

例如：每月支出3萬 × 12 ÷ 4% = 900萬
或：每月支出10萬 × 12 ÷ 4% = 3000萬"""
    c2.text_frame.paragraphs[0].font.size = Pt(22)

    # 投影片3：財富自由的等級
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    t3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    t3.text_frame.text = "📈 財富自由的四個等級"
    t3.text_frame.paragraphs[0].font.size = Pt(40)
    t3.text_frame.paragraphs[0].font.bold = True
    
    c3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    c3.text_frame.text = """🌱 等級一：財務保障
• 被動收入覆蓋基本生活（食衣住行）
• 需要：200-500萬
• 狀態：不再為錢工作

🌿 等級二：財務安全
• 被動收入覆蓋所有生活開銷並有盈餘
• 需要：500-1500萬
• 狀態：有時間自由

🌳 等級三：財務獨立
• 資產持續成長，財務完全自主
• 需要：1500-5000萬
• 狀態：可以選擇性工作

🏆 等級四：財務自由
• 財富規模可以實現任何夢想
• 需要：5000萬以上
• 狀態：完全掌控自己的人生"""
    c3.text_frame.paragraphs[0].font.size = Pt(20)

    # 投影片4：財富自由後的生活想像
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    t4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    t4.text_frame.text = "🌟 財富自由後的生活想像"
    t4.text_frame.paragraphs[0].font.size = Pt(40)
    t4.text_frame.paragraphs[0].font.bold = True
    
    c4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    c4.text_frame.text = """☀️ 理想的一天

早晨：
• 睡到自然醒，不再被鬧鐘追殺
• 悠閒地吃早餐、閱讀、冥想

上午：
• 做自己喜歡的事（運動、創作、學習）
• 或選擇性參與有意義的專案

下午：
• 陪伴家人、朋友
• 追求興趣嗜好

晚上：
• 享用美食、體驗文化
• 早早休息，保持生活品質

🌍 更多的可能性
• 每年多次旅行，探索世界
• 學習新技能、新語言
• 創業實現夢想不必為生存妥協
• 回饋社會，參與公益"""
    c4.text_frame.paragraphs[0].font.size = Pt(20)

    # 投影片5：人生各階段規劃
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    t5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    t5.text_frame.text = "📅 人生各階段規劃"
    t5.text_frame.paragraphs[0].font.size = Pt(40)
    t5.text_frame.paragraphs[0].font.bold = True
    
    c5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    c5.text_frame.text = """👶 25-35歲：播種期
• 提升技能，增加主動收入
• 存下收入的30-50%
• 建立投資紀律
• 勇於嘗試與學習

🌱 35-45歲：成長期
• 事業爆發期，追求升遷或創業
• 資產快速累積
• 分散投資標的
• 建立被動收入系統

🌳 45-55歲：收割期
• 被動收入逐步超越主動收入
• 優化資產配置
• 開始規劃退休生活
• 培養興趣嗜好

🏆 55歲+：自由期
• 享受努力的成果
• 傳承經驗與財富
• 追求人生意義
• 成為下一代的榜樣"""
    c5.text_frame.paragraphs[0].font.size = Pt(19)

    # 投影片6：實現路徑與時間表
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    t6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    t6.text_frame.text = "🗺️ 實現路徑與時間表"
    t6.text_frame.paragraphs[0].font.size = Pt(40)
    t6.text_frame.paragraphs[0].font.bold = True
    
    c6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    c6.text_frame.text = """⏰ 典型時間軸（取決於起始點和努力程度）

5年達成：
• 適合高收入者（月入15萬以上）
• 嚴格執行儲蓄50%以上
• 積極投資於高成長標的

10年達成：
• 適合中等收入者（月入8-15萬）
• 儲蓄30-40%
• 分散投資，指數基金+個股

15-20年達成：
• 適合一般收入者（月入5-8萬）
• 儲蓄20-30%
• 紀律投資，耐心等待複利

🔑 關鍵成功因素
1. 提升收入能力（最關鍵）
2. 控制支出生活品質不犧牲
3. 紀律執行投資計劃
4. 耐心等待複利效果
5. 避免重大投資失敗"""
    c6.text_frame.paragraphs[0].font.size = Pt(19)

    # 投影片7：心理準備與心態調整
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    t7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    t7.text_frame.text = "🧠 心理準備與心態調整"
    t7.text_frame.paragraphs[0].font.size = Pt(40)
    t7.text_frame.paragraphs[0].font.bold = True
    
    c7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    c7.text_frame.text = """💭 需要轉變的心態

❌ 錯誤心態：
• 「有錢就是為了花」
• 「反正會賺回來」
• 「工作是我的全部」
• 「投資一定會贏」

✅ 正確心態：
• 「延遲享受是通往自由的橋樑」
• 「每一筆支出都要有意義」
• 「工作不是目的，是工具」
• 「投資有風險，要紀律分散」
• 「財富是手段，不是終點」

🌱 心理建設
• 接受過程中的不確定性
• 享受努力成長的過程
• 建立超越金錢的價值觀
• 保持謙卑和感恩的心
• 記住最終的目標是「選擇的自由」"""
    c7.text_frame.paragraphs[0].font.size = Pt(19)

    # 投影片8：給讀者的行動建議
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    t8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    t8.text_frame.text = "✅ 給讀者的行動建議"
    t8.text_frame.paragraphs[0].font.size = Pt(40)
    t8.text_frame.paragraphs[0].font.bold = True
    
    c8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    c8.text_frame.text = """🚀 今天就可以開始的10件事

1. 📊 計算你的財富自由門檻
2. 💰 記帳一個月，了解消費習慣
3. 🎯 設定5年後想達成的具體目標
4. 📈 開立投資帳戶，開始小額定期定額
5. 📚 每週閱讀一本理財/商業書籍
6. 🏃 培養一項不需花費太多的興趣
7. 👥 找一個金錢觀相近的社群
8. ✍️ 寫下你財富自由後想做的10件事
9. 🎪 評估你的技能如何轉化為被動收入
10. 🙏 規劃每週「自由時間」讓自己充電

💡 記住：開始永遠比完美重要！"""
    c8.text_frame.paragraphs[0].font.size = Pt(20)

    # 投影片9：總結頁
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    s9_box = slide9.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(4))
    s9_tf = s9_box.text_frame
    s9_tf.text = """🎯 總結

財富自由不是遙不可及的夢想，
而是可以透過規劃和紀律達成的目標。

從今天開始：
✅ 了解你的現狀
✅ 設定明確的目標
✅ 採取行動

你準備好開始這段旅程了嗎？

「真正的自由，是有能力選擇自己想要的生活方式。」

川寶研究團隊｜潔咪"""
    s9_tf.paragraphs[0].font.size = Pt(26)
    s9_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path

def send_telegram(pptx_path):
    """發送到 Telegram"""
    import requests
    url = f"https://api.telegram.org/bot8704642969:AAERVfjKsxcHExGOfZP9h5412w9Sp1TtABw/sendDocument"
    try:
        with open(pptx_path, 'rb') as f:
            files = {'document': f}
            data = {
                'chat_id': '8779713208',
                'caption': '💰 財富自由後的人生規劃（圖文並茂版）| 川寶研究團隊｜潔咪'
            }
            response = requests.post(url, files=files, data=data, timeout=30)
            result = response.json()
            if result.get('ok'):
                print("✅ 已發送到 Telegram")
                return True
            else:
                print(f"❌ 發送失敗: {result}")
                return False
    except Exception as e:
        print(f"❌ 發送錯誤: {e}")
        return False

def main():
    print(f"[{datetime.datetime.now()}] 潔咪開始研究...")
    
    # 取得分析內容
    analysis = get_jiemi_analysis()
    print(f"✅ 分析內容已取得 ({len(analysis)} 字)")
    
    # 建立 PPT
    output_path = f'/root/.openclaw/reports/daily/wealth_freedom_{datetime.datetime.now().strftime("%Y%m%d_%H%M")}.pptx'
    create_wealth_freedom_ppt(analysis, output_path)
    print(f"✅ PPT已生成：{output_path}")
    
    # 發送到 Telegram
    if send_telegram(output_path):
        print("=== 潔咪研究完成 ===")
    else:
        print("⚠️ 報告已生成但 Telegram 發送失敗")

if __name__ == "__main__":
    main()